from pptx.shapes.group import GroupShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import tkinter as tk
from tkinter import filedialog, simpledialog

def check_recursively_for_text(this_set_of_shapes, text_run, slide_number, textbox_counter):
    for shape in this_set_of_shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            textbox_counter = check_recursively_for_text(shape.shapes, text_run, slide_number, textbox_counter)
        else:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    if shape.is_placeholder:
                        placeholder_type = shape.placeholder_format.idx
                        if placeholder_type == 0:
                            category = "Title"
                        elif placeholder_type == 1:
                            category = "Body"
                        else:
                            category = f"Textbox {textbox_counter}"
                            textbox_counter += 1
                    else:
                        category = f"Textbox {textbox_counter}"
                        textbox_counter += 1
                    text_run.append((slide_number, category, text))
            elif shape.has_table:
                category = "Table"
                table_md = []
                for row in shape.table.rows:
                    row_md = []
                    for cell in row.cells:
                        text = " ".join(cell.text.strip().splitlines())
                        row_md.append(text)
                    table_md.append("| " + " | ".join(row_md) + " |")
                if table_md:
                    header = table_md[0]
                    separator = "| " + " | ".join(["---"] * (len(header.split("|")) - 2)) + " |"
                    table_md.insert(1, separator)
                    table_md_str = "\n".join(table_md)
                    text_run.append((slide_number, category, table_md_str))
    return textbox_counter

# Function to process a single file
def process_pptx_file(file_path):
    prs = Presentation(file_path)
    text_run = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        textbox_counter = 1
        textbox_counter = check_recursively_for_text(slide.shapes, text_run, slide_number, textbox_counter)

    def add_table_slide(prs, slide_text, start_index, end_index, part, position):
        slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        title_box = slide.shapes.title

        rows = end_index - start_index + 1
        cols = 2
        left = Inches(0.5)
        top = Inches(0.1)
        width = Inches(8.5)
        height = Inches(0.8)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set column widths
        table.columns[0].width = Inches(2.25)
        table.columns[1].width = Inches(10)

        # Set font size for the first column
        for row in table.rows:
            cell = row.cells[0]
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(16)

        # Set table headers
        table.cell(0, 0).text = f"Title (Part {part})"
        table.cell(0, 1).text = "Body"

        # Set font size for the title row
        for cell in table.rows[0].cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(16)

        # Populate the table with extracted text information
        for i in range(start_index, end_index):
            table.cell(i - start_index + 1, 0).text = slide_text[i]["Title"]
            table.cell(i - start_index + 1, 1).text = slide_text[i]["Body"].strip()

        # Ensure all titles have font size 16
        for i in range(1, len(table.rows)):
            cell = table.cell(i, 0)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(16)

        # Ensure all body text has font size 12
        for i in range(1, len(table.rows)):
            cell = table.cell(i, 1)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)

        # Move the new slide to the specified position
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[-1])
        xml_slides.insert(position, slides[-1])

    # Populate the table with extracted text information
    slide_text = []
    for slide_number, category, text in text_run:
        if category == "Title":
            slide_text.append({"Title": text, "Body": ""})
        else:
            if slide_text:
                slide_text[-1]["Body"] += text + "\n"

    # Add slides with tables
    start_index = 0
    total_lines = 0
    part = 1
    position = 1
    for i in range(len(slide_text)):
        total_lines += slide_text[i]["Body"].count('\n') + 1
        if total_lines > 30:
            add_table_slide(prs, slide_text, start_index, i, part, position)
            start_index = i
            total_lines = slide_text[i]["Body"].count('\n') + 1
            part += 1
            position += 1
    if start_index < len(slide_text):
        add_table_slide(prs, slide_text, start_index, len(slide_text), part, position)

    # Save the modified presentation with "_Modified" appended to the original filename
    base, ext = os.path.splitext(file_path)
    modified_file_path = f"{base}_Modified{ext}"
    prs.save(modified_file_path)

# Function to process all PPTX files in a folder
def process_folder(folder_path, process_subfolders):
    for root, dirs, files in os.walk(folder_path):
        for file in files:
            if file.lower().endswith(".pptx"):
                process_pptx_file(os.path.join(root, file))
        if not process_subfolders:
            break

# Create a Tkinter root window and hide it
root = tk.Tk()
root.withdraw()

# Ask the user to select a file or folder
file_or_folder = filedialog.askopenfilename(title="Select a PPTX file", filetypes=[("PPTX files", "*.pptx")])
if not file_or_folder:
    file_or_folder = filedialog.askdirectory(title="Select a folder containing PPTX files")

# Process the selected file or folder
if os.path.isfile(file_or_folder):
    process_pptx_file(file_or_folder)
elif os.path.isdir(file_or_folder):
    # Ask the user whether to process subfolders
    process_subfolders = tk.messagebox.askyesno("Process Subfolders", "Do you want to process subfolders?")
    process_folder(file_or_folder, process_subfolders)