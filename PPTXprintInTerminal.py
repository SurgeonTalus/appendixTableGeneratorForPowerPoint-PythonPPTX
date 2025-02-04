from pptx.shapes.group import GroupShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from pptx.util import Inches, Pt
import os

def check_recursively_for_text(this_set_of_shapes, text_run, slide_number, textbox_counter):
    for shape in this_set_of_shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            textbox_counter = check_recursively_for_text(shape.shapes, text_run, slide_number, textbox_counter)
        else:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    if shape.is_placeholder:
                        placeholder_type = shape.placeholder_format.idx
                        if placeholder_type == 0:
                            category = "Title"
                        elif placeholder_type == 1:
                            category = "Body"
                        else:
                            category = f"Textbox {textbox_counter}"
                            textbox_counter += 1
                    else:
                        category = f"Textbox {textbox_counter}"
                        textbox_counter += 1
                    text_run.append((slide_number, category, text))
            elif shape.has_table:
                category = "Table"
                table_md = []
                for row in shape.table.rows:
                    row_md = []
                    for cell in row.cells:
                        text = " ".join(cell.text.strip().splitlines())
                        row_md.append(text)
                    table_md.append("| " + " | ".join(row_md) + " |")
                if table_md:
                    header = table_md[0]
                    separator = "| " + " | ".join(["---"] * (len(header.split("|")) - 2)) + " |"
                    table_md.insert(1, separator)
                    table_md_str = "\n".join(table_md)
                    text_run.append((slide_number, category, table_md_str))
    return textbox_counter

def print_text_run(text_run):
    for slide_number, category, text in text_run:
        print(f"Slide {slide_number} - {category}: {text}")

# Specify the file path
file_path = "/Users/sondre/Box Sync/UNDERVISNINGSOPPLEGG/kjetilskolen.com 303 dokumenter Samfunnsfag, Norsk, Engelsk, Spill/gjenfortellingKeynote.pptx"

# Load the presentation
presentation = Presentation(file_path)

text_run = []
for slide_number, slide in enumerate(presentation.slides, start=1):
    check_recursively_for_text(slide.shapes, text_run, slide_number, 0)

# Print the extracted text
print_text_run(text_run)